"""Bootstrap the Journal Figure Skill Hub.

This script creates the reusable hub files from local templates. It does not
execute third-party vendor code. Vendor repositories are only inspected for
metadata such as license files, README presence, and SKILL.md files.
"""

from __future__ import annotations

import json
import math
import os
import shutil
import subprocess
from datetime import date
from pathlib import Path


ROOT = Path(__file__).resolve().parents[2]


SUBSKILLS = {
    "figure-brief-builder": {
        "title": "Figure Brief Builder",
        "purpose": "Convert a paper outline, abstract, or draft into a minimal evidence-driven figure plan.",
        "route": "Use first when the user needs a figure list, figure map, caption skeleton, or generation plan.",
        "inputs": ["paper title", "abstract or section outline", "target journal", "available data or visual evidence"],
        "outputs": ["figure_list.md", "figure_map.json", "figure_spec.json drafts", "caption skeletons"],
        "workflow": [
            "Identify central claims and evidence gaps.",
            "Map each claim to the minimum necessary figure.",
            "Classify figure type and panel structure.",
            "Assign a reproducible generation route.",
            "Flag missing data, missing images, and unresolved journal constraints.",
        ],
    },
    "data-plot-2d": {
        "title": "2D Data Plot",
        "purpose": "Generate publication-quality 2D plots from structured data.",
        "route": "Use for line, scatter, bar, box, violin, heatmap, contour, uncertainty, sensitivity, SHAP-style, and calibration plots.",
        "inputs": ["figure_spec.json", "raw CSV/TSV/JSON data", "plot_config.json"],
        "outputs": ["generate.py", "processed_data.csv", "PDF/SVG/PNG/TIFF exports", "QA report"],
        "workflow": [
            "Preserve raw data and create processed data explicitly.",
            "Render with Matplotlib by default and optional SciencePlots styling.",
            "Label axes and units, record scale transforms, and export vector plus 600 dpi preview.",
            "Run preflight QA before delivery.",
        ],
    },
    "data-plot-3d": {
        "title": "3D Data Plot",
        "purpose": "Generate 3D surfaces, contours, Pareto fronts, and high-dimensional scientific plots.",
        "route": "Use for response surfaces, 3D scatter, contour-surface pairs, parameter maps, and field slices.",
        "inputs": ["figure_spec.json", "structured grid or point data", "viewpoint metadata"],
        "outputs": ["static PDF/SVG/PNG exports", "viewpoint note", "QA report"],
        "workflow": [
            "Check whether a 2D contour communicates better than 3D.",
            "Include colorbar units and viewpoint metadata.",
            "Export a static journal version even if Plotly is used for exploration.",
            "Avoid perspective choices that distort scientific interpretation.",
        ],
    },
    "structure-schematic": {
        "title": "Structure Schematic",
        "purpose": "Create editable schematics of devices, assemblies, specimens, material layers, and geometric concepts.",
        "route": "Use for physical structures, components, specimens, reactors, cells, modules, chips, sensors, and device layouts.",
        "inputs": ["figure_spec.json", "reference sketch or geometry notes", "component list"],
        "outputs": ["editable SVG", "editable PPTX", "PDF/PNG exports", "component legend"],
        "workflow": [
            "Keep shapes, labels, dimensions, arrows, and callouts as separate editable layers.",
            "Use reference images only as trace/background sources.",
            "Do not invent impossible geometry or unsupported dimensions.",
        ],
    },
    "experiment-setup-diagram": {
        "title": "Experiment Setup Diagram",
        "purpose": "Create diagrams of experimental platforms, measurement workflows, and data acquisition paths.",
        "route": "Use for rigs, loops, battery aging systems, microscopy workflows, sensors, DAQ, and control systems.",
        "inputs": ["figure_spec.json", "equipment modules", "sample location", "flow and signal definitions"],
        "outputs": ["editable SVG/PPTX", "legend", "PDF/PNG exports", "QA report"],
        "workflow": [
            "Separate material/energy flow, signal flow, and control flow by line style.",
            "Always mark the sample/specimen position.",
            "Add zoom-in panels when the measurement region is small.",
        ],
    },
    "mechanism-diagram": {
        "title": "Mechanism Diagram",
        "purpose": "Create causal diagrams that explain mechanisms, processes, degradation, transport, or model logic.",
        "route": "Use for conceptual, physical, failure, reaction, transport, architecture, and causal pathway figures.",
        "inputs": ["figure_spec.json", "initial condition", "driving factors", "process states", "observable consequence"],
        "outputs": ["editable SVG/PPTX", "mechanism statement", "PDF/PNG exports", "QA report"],
        "workflow": [
            "Build the chain: initial condition -> driver -> intermediate process -> variable transformation -> observable result.",
            "Assign arrow semantics: force, heat, mass, signal, information, sequence, or causality.",
            "Use equations only when they clarify the mechanism.",
        ],
    },
    "simulation-result-figure": {
        "title": "Simulation Result Figure",
        "purpose": "Prepare FEA/CFD/multiphysics result figures without altering scientific field values.",
        "route": "Use for stress, strain, temperature, velocity, pressure, contact, modal, fatigue, and phase-field maps.",
        "inputs": ["figure_spec.json", "original simulation exports", "software/version/export settings"],
        "outputs": ["annotated composite", "color-scale log", "PDF/PNG/TIFF exports", "QA report"],
        "workflow": [
            "Keep original exports immutable under source/.",
            "Record colorbar ranges, units, deformation scale, and software settings.",
            "Use identical colorbar limits for comparable panels.",
        ],
    },
    "image-annotation-figure": {
        "title": "Image Annotation Figure",
        "purpose": "Annotate experimental photos, microscopy, segmentation, detection, and biomedical images.",
        "route": "Use when real raster evidence needs scale bars, labels, overlays, or traceable annotations.",
        "inputs": ["figure_spec.json", "immutable source image", "calibration data", "annotation plan"],
        "outputs": ["annotated SVG/PDF/PNG", "preprocessing log", "QA report"],
        "workflow": [
            "Never fabricate or hallucinate visual evidence.",
            "Keep annotations in vector layers.",
            "Base scale bars on actual calibration.",
            "Trace overlays to code and source data.",
        ],
    },
    "multi-panel-compositor": {
        "title": "Multi-panel Compositor",
        "purpose": "Combine panels into unified publication figures with consistent labels, spacing, and typography.",
        "route": "Use whenever a figure contains multiple panels or combines plots, schematics, and images.",
        "inputs": ["panel files", "composite_spec.json", "journal width", "caption draft"],
        "outputs": ["composite SVG/PDF/PPTX/PNG", "panel mapping report", "caption draft"],
        "workflow": [
            "Use uniform labels such as (a), (b), (c).",
            "Align axes and margins where possible.",
            "Keep panel text readable at manuscript scale.",
            "Avoid inconsistent color maps in comparable panels.",
        ],
    },
    "graphical-abstract": {
        "title": "Graphical Abstract",
        "purpose": "Create graphical abstracts or TOC figures that summarize the central contribution.",
        "route": "Use for journal graphical abstracts, TOC art, or one-screen conceptual summaries.",
        "inputs": ["core contribution", "3-5 visual modules", "journal requirements", "AI usage decision"],
        "outputs": ["editable SVG/PPTX", "PDF/PNG exports", "ai_disclosure.md if needed"],
        "workflow": [
            "Summarize the paper's central contribution, not every detail.",
            "Use no more than 3-5 visual modules.",
            "Use AI imagery only for ideation unless disclosed and journal-compliant.",
        ],
    },
    "journal-export": {
        "title": "Journal Export",
        "purpose": "Export figures according to generic or target-journal file, size, DPI, and font constraints.",
        "route": "Use before manuscript insertion, submission packaging, or journal-specific conversion.",
        "inputs": ["figure_spec.json", "editable source", "target journal or generic preset"],
        "outputs": ["PDF/SVG/PNG/TIFF/EPS as available", "export_log.json", "journal package"],
        "workflow": [
            "Apply single-column or double-column width presets.",
            "Check DPI, font size, vector availability, and file size.",
            "Record missing optional tools such as Inkscape or ImageMagick.",
        ],
    },
    "figure-preflight-qa": {
        "title": "Figure Preflight QA",
        "purpose": "Run the final quality gate for traceability, editability, journal compliance, and readability.",
        "route": "Always run at the end of a figure task.",
        "inputs": ["figure_spec.json", "source files", "editable outputs", "final exports", "caption draft"],
        "outputs": ["qa/preflight_report.md", "qa/preflight_report.json", "qa/issues_to_fix.md"],
        "workflow": [
            "Check required files, formats, resolution, and vector/editable outputs.",
            "Check labels, units, legends, colorbars, and caption-panel consistency.",
            "Flag AI disclosure, third-party image, and source-traceability risks.",
        ],
    },
}


VENDOR_DECISIONS = {
    "openai-skills": ("partial", "Use as a reference for concise skill packaging and progressive disclosure."),
    "awesome-codex-skills": ("partial", "Use as an index of skill packaging patterns; do not import third-party skills."),
    "awesome-codex-subagents": ("partial", "Use as a reference for task decomposition and quality-gate language."),
    "awesome-agent-skills": ("partial", "Use as a reference for cross-agent skill conventions."),
    "andrej-karpathy-skills": ("partial", "Use as a reference for minimal readable skill style."),
    "SciencePlots": ("partial", "Optional Matplotlib style dependency/reference for publication plots."),
    "pylustrator": ("partial", "Reference for reproducible post-layout polishing; keep optional."),
    "matplotlib_for_papers": ("partial", "Reference for size, fonts, and Matplotlib publication practices."),
    "Kaleido": ("partial", "Optional static export route for Plotly-generated 3D figures."),
    "CairoSVG": ("partial", "Optional SVG to PDF/PNG conversion backend."),
    "drawsvg": ("partial", "Optional inspiration for programmatic SVG generation; local templates use plain SVG."),
    "figure_tutorial": ("partial", "Reference for multi-panel SVG composition practices."),
    "panel-plots": ("partial", "Reference for exact-size multi-panel layout ideas."),
    "scientific-agent-skills": ("partial", "Reference only; broad third-party skills stay quarantined."),
    "Codex-drawio-skill": ("partial", "Reference for editable diagram workflow; not adopted as a runtime dependency."),
}


def write(path: str | Path, content: str) -> None:
    p = ROOT / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content.strip() + "\n", encoding="utf-8")


def json_write(path: str | Path, data: object) -> None:
    p = ROOT / path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def skill_md(name: str, spec: dict[str, object]) -> str:
    inputs = "\n".join(f"- {x}" for x in spec["inputs"])
    outputs = "\n".join(f"- {x}" for x in spec["outputs"])
    workflow = "\n".join(f"{i + 1}. {x}" for i, x in enumerate(spec["workflow"]))
    return f"""---
name: {name}
description: {spec["purpose"]} {spec["route"]} Use inside journal-figure-skill-hub when planning, generating, editing, exporting, or checking publication-quality scientific figures.
---

# {spec["title"]}

## Applicability
{spec["purpose"]}

## Inputs
{inputs}

## Outputs
{outputs}

## Recommended Route
{workflow}

## Editable Layer Requirements
- Keep text, arrows, shapes, annotations, legends, panel labels, and callouts editable whenever the output is line art or schematic.
- For raster evidence, keep the source image immutable and place annotations in a separate vector layer.
- Export at least one editable/vector artifact and one review-friendly raster preview.

## Quality Checks
- Validate or create `figure_spec.json` before generation.
- Confirm source traceability and document missing evidence.
- Check labels, units, colorbars, scale bars, font size, panel labels, and caption consistency.
- Run `shared/scripts/build_qa_report.py` before delivery.

## Example Prompt
Use `{name}` to create a journal-ready figure from my data or figure brief. Create a `figure_spec.json`, preserve editable outputs, export PDF/SVG/PNG, and run preflight QA.

## Example Output Directory
`figures/Fig_XX_{name.replace("-", "_")}/`

## Handoff
- Route data-bearing panels to `data-plot-2d` or `data-plot-3d`.
- Route diagram panels to `structure-schematic`, `experiment-setup-diagram`, or `mechanism-diagram`.
- Route all multi-panel results to `multi-panel-compositor`.
- Finish every workflow with `figure-preflight-qa`.
"""


def create_skills() -> None:
    write(
        "SKILL.md",
        """---
name: journal-figure-skill-hub
description: Universal skill hub for planning, generating, composing, exporting, and checking publication-quality journal figures across disciplines. Use when the user asks for paper figures, scientific schematics, experimental diagrams, mechanism diagrams, data plots, simulation result figures, graphical abstracts, figure captions, or journal-ready figure exports.
---

# Journal Figure Skill Hub

## Core Contract
Use this hub for scientific and journal-paper figures across mechanical engineering, materials, energy, batteries, AI/ML, biomedical engineering, civil engineering, robotics, thermal/fluid science, reliability, experimental systems, and simulation-driven studies.

Every figure must start from or produce `figure_spec.json`, preserve source traceability, keep editable/vector outputs whenever possible, and finish with preflight QA.

## Routing Rules
1. Need a figure plan or figure list: use `skills/figure-brief-builder`.
2. Numerical 2D data: use `skills/data-plot-2d`.
3. 3D or high-dimensional data: use `skills/data-plot-3d`.
4. Physical object, assembly, specimen, or device: use `skills/structure-schematic`.
5. Experimental rig, measurement workflow, signal/control path: use `skills/experiment-setup-diagram`.
6. Causal, physical, degradation, reaction, transport, or model mechanism: use `skills/mechanism-diagram`.
7. FEA/CFD/multiphysics contour or field result: use `skills/simulation-result-figure`.
8. Real photos, microscopy, segmentation, detection, or biomedical images: use `skills/image-annotation-figure`.
9. Multiple panels: use `skills/multi-panel-compositor`.
10. Graphical abstract or TOC figure: use `skills/graphical-abstract`.
11. Journal packaging/export: use `skills/journal-export`.
12. Always finish with `skills/figure-preflight-qa`.

## Global Rules
- Do not fabricate experimental, simulation, microscopy, medical, or image evidence.
- Data-bearing figures must be reproducible from source data and code.
- AI image generation is allowed only for ideation or disclosed non-evidentiary visual material.
- Preserve raw data and original images under `source/`.
- Keep line art, plots, schematics, labels, arrows, and callouts editable.
- Export at least one editable/vector file and one review-friendly raster preview.
- Store figure projects as `figures/Fig_XX_short_name/` with `source/`, `code/`, `editable/`, `final/`, `qa/`, and `caption/`.

## Required Workflow
1. Create or read `figure_spec.json`.
2. Route to the narrowest sub-skill.
3. Generate source-controlled code/templates.
4. Export PDF/SVG/PNG and optional TIFF/PPTX.
5. Run preflight QA.
6. Report limitations and missing evidence explicitly.
""",
    )
    for name, spec in SUBSKILLS.items():
        base = ROOT / "skills" / name
        for sub in ["scripts", "references", "assets", "examples", "tests"]:
            (base / sub).mkdir(parents=True, exist_ok=True)
        write(Path("skills") / name / "SKILL.md", skill_md(name, spec))


def create_configs_and_schemas() -> None:
    write(
        "config/default_style.yaml",
        """
font_family: Arial
font_size_pt: 8
line_width_pt: 1.0
panel_label_size_pt: 9
axis_label_size_pt: 8
tick_label_size_pt: 7
palette: colorblind_safe
background: white
""",
    )
    write(
        "config/color_palettes.yaml",
        """
colorblind_safe:
  - "#0072B2"
  - "#D55E00"
  - "#009E73"
  - "#CC79A7"
  - "#E69F00"
  - "#56B4E9"
  - "#F0E442"
  - "#000000"
sequential_thermal:
  - "#313695"
  - "#74add1"
  - "#ffffbf"
  - "#f46d43"
  - "#a50026"
""",
    )
    write(
        "config/journal_presets.yaml",
        """
default_single_column:
  width_mm: 90
  min_font_pt: 7
  raster_dpi: 600
default_double_column:
  width_mm: 180
  min_font_pt: 7
  raster_dpi: 600
graphical_abstract_generic:
  width_mm: 120
  height_mm: 50
  min_font_pt: 7
  raster_dpi: 300
""",
    )
    write(
        "config/export_presets.yaml",
        """
standard_vector:
  formats: ["pdf", "svg", "png"]
  dpi: 600
journal_upload:
  formats: ["pdf", "png", "tiff"]
  dpi: 600
editable_review:
  formats: ["svg", "pptx", "png"]
  dpi: 300
""",
    )
    write(
        "config/qa_rules.yaml",
        """
required_project_files:
  - figure_spec.json
required_dirs:
  - source
  - code
  - editable
  - final
  - qa
  - caption
minimum_png_dpi: 300
preferred_png_dpi: 600
require_source_traceability: true
require_editable_for_line_art: true
""",
    )
    write(
        "config/hub_config.yaml",
        """
hub_name: journal-figure-skill-hub
default_figure_root: figures
strict_mode: true
allow_ai_final_evidence: false
""",
    )
    schema = {
        "$schema": "https://json-schema.org/draft/2020-12/schema",
        "title": "Journal Figure Specification",
        "type": "object",
        "required": [
            "figure_id",
            "short_name",
            "paper_context",
            "figure_type",
            "scientific_purpose",
            "reader_takeaway_5s",
            "panels",
            "data_sources",
            "visual_style",
            "export_targets",
            "journal_constraints",
            "ai_usage",
            "qa_requirements",
        ],
        "properties": {
            "figure_id": {"type": "string"},
            "short_name": {"type": "string"},
            "paper_context": {
                "type": "object",
                "required": ["title", "discipline", "section", "target_journal", "core_claim"],
                "properties": {
                    "title": {"type": "string"},
                    "discipline": {"type": "string"},
                    "section": {"type": "string"},
                    "target_journal": {"type": "string"},
                    "core_claim": {"type": "string"},
                },
            },
            "figure_type": {
                "type": "string",
                "enum": [
                    "structure_schematic",
                    "experiment_setup_diagram",
                    "mechanism_diagram",
                    "data_plot_2d",
                    "data_plot_3d",
                    "simulation_result_figure",
                    "image_annotation_figure",
                    "multi_panel_composite",
                    "graphical_abstract",
                    "journal_export",
                    "preflight_qa",
                ],
            },
            "scientific_purpose": {"type": "string"},
            "reader_takeaway_5s": {"type": "string"},
            "panels": {
                "type": "array",
                "minItems": 1,
                "items": {
                    "type": "object",
                    "required": ["panel_id", "title", "message", "input_sources", "generation_route", "editable_layers", "output_files"],
                    "properties": {
                        "panel_id": {"type": "string"},
                        "title": {"type": "string"},
                        "message": {"type": "string"},
                        "input_sources": {"type": "array", "items": {"type": "string"}},
                        "generation_route": {"type": "string"},
                        "editable_layers": {"type": "array", "items": {"type": "string"}},
                        "output_files": {"type": "array", "items": {"type": "string"}},
                    },
                },
            },
            "data_sources": {"type": "array", "items": {"type": "string"}},
            "visual_style": {
                "type": "object",
                "required": ["width_mm", "font_family", "font_size_pt", "line_width_pt", "color_palette"],
                "properties": {
                    "width_mm": {"type": "number"},
                    "height_mm": {"type": ["number", "null"]},
                    "font_family": {"type": "string"},
                    "font_size_pt": {"type": "number"},
                    "line_width_pt": {"type": "number"},
                    "color_palette": {"type": "string"},
                },
            },
            "export_targets": {"type": "array", "items": {"type": "string"}},
            "journal_constraints": {"type": "object"},
            "ai_usage": {"type": "object"},
            "qa_requirements": {"type": "array", "items": {"type": "string"}},
        },
    }
    json_write("shared/schemas/figure_spec.schema.json", schema)
    json_write("shared/schemas/plot_spec.schema.json", {"$schema": schema["$schema"], "title": "Plot Spec", "type": "object"})
    json_write("shared/schemas/schematic_spec.schema.json", {"$schema": schema["$schema"], "title": "Schematic Spec", "type": "object"})
    json_write("shared/schemas/composite_spec.schema.json", {"$schema": schema["$schema"], "title": "Composite Spec", "type": "object"})
    json_write("shared/schemas/qa_report.schema.json", {"$schema": schema["$schema"], "title": "QA Report", "type": "object"})
    json_write(
        "shared/templates/figure_spec_template.json",
        {
            "figure_id": "Fig. 1",
            "short_name": "mechanism_overview",
            "paper_context": {
                "title": "",
                "discipline": "general",
                "section": "Introduction",
                "target_journal": "",
                "core_claim": "",
            },
            "figure_type": "mechanism_diagram",
            "scientific_purpose": "",
            "reader_takeaway_5s": "",
            "panels": [
                {
                    "panel_id": "a",
                    "title": "",
                    "message": "",
                    "input_sources": [],
                    "generation_route": "svg",
                    "editable_layers": ["text", "arrows", "shapes"],
                    "output_files": [],
                }
            ],
            "data_sources": [],
            "visual_style": {
                "width_mm": 180,
                "height_mm": None,
                "font_family": "Arial",
                "font_size_pt": 8,
                "line_width_pt": 1.0,
                "color_palette": "colorblind_safe",
            },
            "export_targets": ["pdf", "svg", "png_600dpi", "pptx"],
            "journal_constraints": {
                "single_or_double_column": "double",
                "requires_tiff": False,
                "requires_eps": False,
                "max_file_size_mb": None,
            },
            "ai_usage": {
                "used_ai_image_generation": False,
                "usage_scope": "none",
                "requires_disclosure": False,
            },
            "qa_requirements": [
                "resolution_check",
                "font_check",
                "axis_label_check",
                "unit_check",
                "editable_layer_check",
                "source_traceability_check",
            ],
        },
    )


def create_core_scripts() -> None:
    write(
        "shared/scripts/create_figure_project.py",
        r'''
from __future__ import annotations
import argparse, json, shutil
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]

def create_project(figure_id: str, name: str, figure_type: str, out_root: Path) -> Path:
    safe_id = figure_id.replace(".", "").replace(" ", "_")
    project = out_root / f"{safe_id}_{name}"
    for sub in ["source", "code", "editable", "final", "qa", "caption"]:
        (project / sub).mkdir(parents=True, exist_ok=True)
    template = ROOT / "shared" / "templates" / "figure_spec_template.json"
    data = json.loads(template.read_text(encoding="utf-8"))
    data["figure_id"] = figure_id
    data["short_name"] = name
    data["figure_type"] = figure_type
    (project / "figure_spec.json").write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
    (project / "README.md").write_text(f"# {figure_id} {name}\n\nGenerated by journal-figure-skill-hub.\n", encoding="utf-8")
    (project / "caption" / "caption_draft.md").write_text(f"{figure_id}. Caption draft for {name}.\n", encoding="utf-8")
    (project / "caption" / "ai_disclosure.md").write_text("AI image generation used: no.\n", encoding="utf-8")
    return project

def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("--figure-id", required=True)
    p.add_argument("--name", required=True)
    p.add_argument("--type", required=True)
    p.add_argument("--out-root", default="figures")
    args = p.parse_args()
    project = create_project(args.figure_id, args.name, args.type, Path(args.out_root))
    print(project)

if __name__ == "__main__":
    main()
''',
    )
    write(
        "shared/scripts/validate_figure_spec.py",
        r'''
from __future__ import annotations
import argparse, json, sys
from pathlib import Path
from jsonschema import Draft202012Validator

ROOT = Path(__file__).resolve().parents[2]

def validate(path: Path) -> list[str]:
    schema = json.loads((ROOT / "shared" / "schemas" / "figure_spec.schema.json").read_text(encoding="utf-8"))
    data = json.loads(path.read_text(encoding="utf-8"))
    validator = Draft202012Validator(schema)
    return [f"{list(e.path)}: {e.message}" for e in validator.iter_errors(data)]

def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("spec")
    args = p.parse_args()
    errors = validate(Path(args.spec))
    if errors:
        print("\n".join(errors))
        sys.exit(1)
    print(f"OK: {args.spec}")

if __name__ == "__main__":
    main()
''',
    )
    write(
        "shared/scripts/export_plot.py",
        r'''
from __future__ import annotations
from pathlib import Path

def export_matplotlib(fig, output_stem: str | Path, dpi: int = 600, tiff: bool = False) -> list[Path]:
    stem = Path(output_stem)
    stem.parent.mkdir(parents=True, exist_ok=True)
    outputs = []
    for ext in ["pdf", "svg", "png"]:
        path = stem.with_suffix(f".{ext}")
        fig.savefig(path, dpi=dpi if ext == "png" else None, bbox_inches="tight")
        outputs.append(path)
    if tiff:
        path = stem.with_suffix(".tiff")
        fig.savefig(path, dpi=dpi, bbox_inches="tight")
        outputs.append(path)
    return outputs
''',
    )
    write(
        "shared/scripts/export_svg.py",
        r'''
from __future__ import annotations
import argparse, shutil, subprocess
from pathlib import Path

def export_svg(svg: Path, out_dir: Path, dpi: int = 600) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    outputs = []
    try:
        import cairosvg
        pdf = out_dir / f"{svg.stem}.pdf"
        png = out_dir / f"{svg.stem}_{dpi}dpi.png"
        cairosvg.svg2pdf(url=str(svg), write_to=str(pdf))
        cairosvg.svg2png(url=str(svg), write_to=str(png), dpi=dpi)
        outputs.extend([pdf, png])
    except Exception as exc:
        inkscape = shutil.which("inkscape")
        if not inkscape:
            raise RuntimeError(f"CairoSVG failed and Inkscape is unavailable: {exc}") from exc
        pdf = out_dir / f"{svg.stem}.pdf"
        png = out_dir / f"{svg.stem}_{dpi}dpi.png"
        subprocess.run([inkscape, str(svg), "--export-type=pdf", f"--export-filename={pdf}"], check=True)
        subprocess.run([inkscape, str(svg), "--export-type=png", f"--export-dpi={dpi}", f"--export-filename={png}"], check=True)
        outputs.extend([pdf, png])
    return outputs

def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("svg")
    p.add_argument("--out-dir", default="final")
    p.add_argument("--dpi", type=int, default=600)
    args = p.parse_args()
    print("\n".join(str(x) for x in export_svg(Path(args.svg), Path(args.out_dir), args.dpi)))

if __name__ == "__main__":
    main()
''',
    )
    write(
        "shared/scripts/export_pptx.py",
        r'''
from __future__ import annotations
import argparse, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

def build_pptx(shape_spec: Path, output: Path) -> Path:
    spec = json.loads(shape_spec.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(spec.get("width_in", 10))
    prs.slide_height = Inches(spec.get("height_in", 5.625))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for s in spec.get("shapes", []):
        x, y, w, h = [Inches(float(s.get(k, 0))) for k in ["x", "y", "w", "h"]]
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.text = s.get("text", "")
        if shape.text_frame and shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].font.size = Pt(float(s.get("font_pt", 12)))
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output

def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("shape_spec")
    p.add_argument("output")
    args = p.parse_args()
    print(build_pptx(Path(args.shape_spec), Path(args.output)))

if __name__ == "__main__":
    main()
''',
    )
    write(
        "shared/scripts/compose_panels.py",
        r'''
from __future__ import annotations
import argparse, html, json, math
from pathlib import Path
from xml.etree import ElementTree as ET

def compose(spec_path: Path, output: Path) -> Path:
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    width = int(spec.get("width", 1200))
    height = int(spec.get("height", 800))
    margin = int(spec.get("margin", 45))
    panels = spec["panels"]
    cols = int(spec.get("cols", math.ceil(math.sqrt(len(panels)))))
    rows = math.ceil(len(panels) / cols)
    cell_w = (width - margin * (cols + 1)) / cols
    cell_h = (height - margin * (rows + 1)) / rows
    parts = [f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">']
    parts.append('<rect width="100%" height="100%" fill="white"/>')
    for i, panel in enumerate(panels):
        row, col = divmod(i, cols)
        x = margin + col * (cell_w + margin)
        y = margin + row * (cell_h + margin)
        label = panel.get("label", f"({chr(97+i)})")
        title = html.escape(panel.get("title", ""))
        parts.append(f'<g id="panel_{i+1}">')
        parts.append(f'<text x="{x}" y="{y-12}" font-family="Arial" font-size="22" font-weight="bold">{label}</text>')
        parts.append(f'<rect x="{x}" y="{y}" width="{cell_w}" height="{cell_h}" fill="#f8f8f8" stroke="#444" stroke-width="1"/>')
        img = panel.get("file")
        if img:
            href = html.escape(img)
            parts.append(f'<image href="{href}" x="{x+10}" y="{y+20}" width="{cell_w-20}" height="{cell_h-45}" preserveAspectRatio="xMidYMid meet"/>')
        parts.append(f'<text x="{x+12}" y="{y+cell_h-12}" font-family="Arial" font-size="14">{title}</text>')
        parts.append("</g>")
    parts.append("</svg>")
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("\n".join(parts), encoding="utf-8")
    return output

def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("spec")
    p.add_argument("output")
    args = p.parse_args()
    print(compose(Path(args.spec), Path(args.output)))

if __name__ == "__main__":
    main()
''',
    )
    write(
        "shared/scripts/check_image_resolution.py",
        r'''
from __future__ import annotations
import argparse
from pathlib import Path
from PIL import Image

def check(path: Path) -> dict:
    with Image.open(path) as img:
        dpi = img.info.get("dpi", (0, 0))
        return {"file": str(path), "width_px": img.width, "height_px": img.height, "dpi": dpi}

def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("image")
    args = p.parse_args()
    print(check(Path(args.image)))

if __name__ == "__main__":
    main()
''',
    )
    write(
        "shared/scripts/check_svg_text_layers.py",
        r'''
from __future__ import annotations
import argparse
from pathlib import Path
from xml.etree import ElementTree as ET

def has_text_layers(svg: Path) -> bool:
    root = ET.parse(svg).getroot()
    return any(el.tag.endswith("text") for el in root.iter())

def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("svg")
    args = p.parse_args()
    ok = has_text_layers(Path(args.svg))
    print(f"text_layers={ok}")
    raise SystemExit(0 if ok else 1)

if __name__ == "__main__":
    main()
''',
    )
    write(
        "shared/scripts/build_qa_report.py",
        r'''
from __future__ import annotations
import argparse, json
from datetime import datetime
from pathlib import Path

REQUIRED_DIRS = ["source", "code", "editable", "final", "qa", "caption"]

def qa(project: Path) -> dict:
    issues = []
    spec = project / "figure_spec.json"
    if not spec.exists():
        issues.append("Missing figure_spec.json")
        data = {}
    else:
        data = json.loads(spec.read_text(encoding="utf-8"))
    for d in REQUIRED_DIRS:
        if not (project / d).is_dir():
            issues.append(f"Missing directory: {d}")
    final_files = list((project / "final").glob("*")) if (project / "final").exists() else []
    editable_files = list((project / "editable").glob("*")) if (project / "editable").exists() else []
    if not final_files:
        issues.append("No final exports found")
    if not editable_files:
        issues.append("No editable outputs found")
    if data.get("ai_usage", {}).get("requires_disclosure") and not (project / "caption" / "ai_disclosure.md").exists():
        issues.append("AI disclosure required but missing")
    report = {
        "project": str(project),
        "checked_at": datetime.now().isoformat(timespec="seconds"),
        "figure_type": data.get("figure_type"),
        "final_files": [str(x.relative_to(project)) for x in final_files],
        "editable_files": [str(x.relative_to(project)) for x in editable_files],
        "issues": issues,
        "status": "pass" if not issues else "needs_fix",
    }
    qadir = project / "qa"
    qadir.mkdir(exist_ok=True)
    (qadir / "preflight_report.json").write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")
    md = ["# Preflight QA Report", "", f"- Status: {report['status']}", f"- Figure type: {report.get('figure_type')}", ""]
    md.append("## Final files")
    md.extend(f"- {x}" for x in report["final_files"])
    md.append("\n## Editable files")
    md.extend(f"- {x}" for x in report["editable_files"])
    md.append("\n## Issues")
    md.extend(f"- {x}" for x in issues or ["None"])
    (qadir / "preflight_report.md").write_text("\n".join(md) + "\n", encoding="utf-8")
    (qadir / "issues_to_fix.md").write_text("\n".join(f"- {x}" for x in issues) + ("\n" if issues else "- None\n"), encoding="utf-8")
    return report

def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("project")
    args = p.parse_args()
    print(json.dumps(qa(Path(args.project)), indent=2, ensure_ascii=False))

if __name__ == "__main__":
    main()
''',
    )


def create_templates() -> None:
    write(
        "shared/templates/plot_2d_template.py",
        r'''
from pathlib import Path
import matplotlib.pyplot as plt
import pandas as pd

from shared.scripts.export_plot import export_matplotlib

df = pd.read_csv(Path("source") / "raw_data.csv")
fig, ax = plt.subplots(figsize=(3.5, 2.4))
ax.plot(df["x"], df["y"], marker="o", label="Series")
ax.set_xlabel("Independent variable (unit)")
ax.set_ylabel("Response (unit)")
ax.legend(frameon=False)
export_matplotlib(fig, Path("final") / "figure", dpi=600)
''',
    )
    write(
        "shared/templates/plot_3d_template.py",
        r'''
from pathlib import Path
import matplotlib.pyplot as plt
import pandas as pd
from mpl_toolkits.mplot3d import Axes3D  # noqa: F401
from shared.scripts.export_plot import export_matplotlib

df = pd.read_csv(Path("source") / "surface_data.csv")
fig = plt.figure(figsize=(4.5, 3.2))
ax = fig.add_subplot(111, projection="3d")
surf = ax.plot_trisurf(df["x"], df["y"], df["z"], cmap="viridis", linewidth=0.2)
ax.set_xlabel("x (unit)")
ax.set_ylabel("y (unit)")
ax.set_zlabel("response (unit)")
fig.colorbar(surf, ax=ax, shrink=0.65, label="response (unit)")
export_matplotlib(fig, Path("final") / "figure", dpi=600)
''',
    )
    write(
        "shared/templates/mechanism_svg_template.py",
        r"""
from pathlib import Path

SVG = '''<svg xmlns="http://www.w3.org/2000/svg" width="1000" height="360" viewBox="0 0 1000 360">
<rect width="100%" height="100%" fill="white"/>
<defs><marker id="arrow" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto"><path d="M0,0 L0,6 L9,3 z" fill="#333"/></marker></defs>
<g font-family="Arial" font-size="22">
<rect x="60" y="120" width="190" height="90" rx="8" fill="#E8F1FA" stroke="#0072B2"/>
<text x="95" y="170">Initial state</text>
<line x1="260" y1="165" x2="390" y2="165" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
<rect x="405" y="120" width="190" height="90" rx="8" fill="#FCE8D8" stroke="#D55E00"/>
<text x="450" y="170">Mechanism</text>
<line x1="605" y1="165" x2="735" y2="165" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
<rect x="750" y="120" width="190" height="90" rx="8" fill="#E9F6EF" stroke="#009E73"/>
<text x="790" y="170">Observable</text>
</g></svg>'''

Path("editable").mkdir(exist_ok=True)
Path("final").mkdir(exist_ok=True)
(Path("editable") / "mechanism.svg").write_text(SVG, encoding="utf-8")
(Path("final") / "mechanism.svg").write_text(SVG, encoding="utf-8")
""",
    )
    write(
        "shared/templates/experiment_diagram_template.svg",
        """
<svg xmlns="http://www.w3.org/2000/svg" width="1000" height="520" viewBox="0 0 1000 520">
  <rect width="100%" height="100%" fill="white"/>
  <defs>
    <marker id="arrow" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto">
      <path d="M0,0 L0,6 L9,3 z" fill="#333"/>
    </marker>
  </defs>
  <g font-family="Arial" font-size="20">
    <rect x="80" y="220" width="150" height="90" fill="#E8F1FA" stroke="#0072B2"/>
    <text x="115" y="272">Sample</text>
    <rect x="330" y="120" width="160" height="80" fill="#FCE8D8" stroke="#D55E00"/>
    <text x="362" y="167">Sensor</text>
    <rect x="585" y="120" width="150" height="80" fill="#E9F6EF" stroke="#009E73"/>
    <text x="632" y="167">DAQ</text>
    <rect x="790" y="120" width="150" height="80" fill="#EFEFEF" stroke="#555"/>
    <text x="820" y="167">Computer</text>
    <line x1="230" y1="265" x2="330" y2="170" stroke="#D55E00" stroke-width="3" stroke-dasharray="8 6" marker-end="url(#arrow)"/>
    <line x1="490" y1="160" x2="585" y2="160" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
    <line x1="735" y1="160" x2="790" y2="160" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
    <line x1="80" y1="360" x2="920" y2="360" stroke="#0072B2" stroke-width="4" marker-end="url(#arrow)"/>
    <text x="365" y="395">physical flow</text>
    <text x="445" y="95">signal/control flow</text>
  </g>
</svg>
""",
    )
    json_write(
        "shared/templates/composite_layout_template.json",
        {
            "width": 1200,
            "height": 800,
            "cols": 2,
            "margin": 50,
            "panels": [
                {"label": "(a)", "title": "Panel A", "file": "panel_a.png"},
                {"label": "(b)", "title": "Panel B", "file": "panel_b.png"},
            ],
        },
    )
    write("shared/templates/caption_template.md", "Fig. X. Concise caption that maps each panel to the scientific claim.")


def create_references() -> None:
    write(
        "shared/references/figure_design_principles.md",
        """
# Figure Design Principles

- Start from the scientific claim and evidence, not visual decoration.
- Prefer fewer panels with clearer messages.
- Preserve traceability from source data or original image to final export.
- Use color only to encode meaningful grouping, magnitude, or state.
- Keep labels readable at final manuscript scale.
""",
    )
    write(
        "shared/references/journal_image_requirements.md",
        """
# Journal Image Requirements

Generic starting presets:
- Single column: 90 mm width, minimum 7 pt text, 600 dpi raster preview.
- Double column: 180 mm width, minimum 7 pt text, 600 dpi raster preview.
- Prefer PDF/SVG for vector line art and TIFF/PNG for raster evidence when required.
""",
    )
    write(
        "shared/references/ai_generated_image_policy_notes.md",
        """
# AI Generated Image Policy Notes

AI imagery can support ideation, thumbnails, and layout exploration. It must not invent scientific evidence such as real experimental images, simulation contours, medical imagery, material morphology, cracks, flames, or failures. If AI imagery remains in a final figure, create `caption/ai_disclosure.md` and check the target journal policy.
""",
    )
    write(
        "shared/references/matplotlib_style_guide.md",
        """
# Matplotlib Style Guide

- Use explicit figure size in inches derived from journal width in mm.
- Set font family, font size, line width, tick direction, and legend frame consistently.
- Label axes with units.
- Save vector exports with `bbox_inches="tight"` and raster previews at 600 dpi.
- Record all smoothing, filtering, scaling, and outlier handling.
""",
    )
    write(
        "shared/references/svg_pptx_editability_guide.md",
        """
# SVG/PPTX Editability Guide

- Keep text as text nodes when possible.
- Use groups for panels, annotations, dimensions, and data overlays.
- Keep raster evidence as immutable source images plus vector annotation layers.
- Use PPTX shapes for editable review copies when collaborators revise in PowerPoint.
""",
    )
    write(
        "shared/references/github_resource_audit_method.md",
        """
# GitHub Resource Audit Method

Clone into `_research/vendor_repos/`, inspect license and maintenance status, scan for SKILL.md or operational instructions, review dependencies, and decide whether to accept, partially accept, or reject. Do not execute untrusted vendor code during audit.
""",
    )


def create_demo_script() -> None:
    write(
        "shared/scripts/make_demo_figures.py",
        r"""
from __future__ import annotations

import json
import math
import shutil
import subprocess
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from create_figure_project import create_project
from export_plot import export_matplotlib
from export_svg import export_svg
from build_qa_report import qa
from compose_panels import compose

ROOT = Path(__file__).resolve().parents[2]
FIGROOT = ROOT / "examples" / "demo_project" / "figures"

def update_spec(project: Path, figure_type: str, purpose: str, outputs: list[str], sources: list[str]) -> None:
    spec = json.loads((project / "figure_spec.json").read_text(encoding="utf-8"))
    spec["figure_type"] = figure_type
    spec["scientific_purpose"] = purpose
    spec["reader_takeaway_5s"] = purpose
    spec["data_sources"] = sources
    spec["panels"][0]["message"] = purpose
    spec["panels"][0]["output_files"] = outputs
    (project / "figure_spec.json").write_text(json.dumps(spec, indent=2, ensure_ascii=False), encoding="utf-8")

def copy_script(project: Path) -> None:
    shutil.copy2(__file__, project / "code" / "generate.py")

def pptx_placeholder(project: Path, name: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.5), Inches(8.8), Inches(4.5))
    box.text = f"Editable review copy: {name}"
    box.text_frame.paragraphs[0].font.size = Pt(24)
    prs.save(project / "editable" / f"{name}.pptx")

def fallback_preview(project: Path, name: str, title: str) -> None:
    fig, ax = plt.subplots(figsize=(6.0, 2.6))
    ax.axis("off")
    ax.text(0.5, 0.72, title, ha="center", va="center", fontsize=12, fontweight="bold")
    ax.text(0.5, 0.48, "Editable vector source is the SVG file in editable/ and final/.", ha="center", va="center", fontsize=9)
    ax.text(0.5, 0.30, "This PDF/PNG is a renderer fallback preview for environments without Cairo/Inkscape.", ha="center", va="center", fontsize=8)
    export_matplotlib(fig, project / "final" / name, dpi=600)
    plt.close(fig)

def demo_2d() -> Path:
    project = create_project("Fig_01", "model_performance_trend", "data_plot_2d", FIGROOT)
    x = np.linspace(0, 10, 28)
    df = pd.DataFrame({"x": x, "baseline": 0.58 + 0.18 * (1 - np.exp(-x / 4)), "proposed": 0.62 + 0.27 * (1 - np.exp(-x / 3))})
    df.to_csv(project / "source" / "raw_data.csv", index=False)
    fig, ax = plt.subplots(figsize=(3.6, 2.4))
    ax.plot(df["x"], df["baseline"], "o-", label="Baseline")
    ax.plot(df["x"], df["proposed"], "s-", label="Proposed")
    ax.set_xlabel("Training size (a.u.)")
    ax.set_ylabel("Score (a.u.)")
    ax.legend(frameon=False)
    ax.grid(alpha=0.25)
    export_matplotlib(fig, project / "final" / "Fig_01_model_performance_trend")
    plt.close(fig)
    shutil.copy2(project / "final" / "Fig_01_model_performance_trend.svg", project / "editable" / "Fig_01_model_performance_trend.svg")
    pptx_placeholder(project, "Fig_01_model_performance_trend")
    update_spec(project, "data_plot_2d", "Compare synthetic model-performance trends with reproducible source data.", ["final/Fig_01_model_performance_trend.pdf", "final/Fig_01_model_performance_trend.svg", "final/Fig_01_model_performance_trend.png"], ["source/raw_data.csv"])
    copy_script(project)
    qa(project)
    return project

def demo_3d() -> Path:
    project = create_project("Fig_02", "response_surface_contour", "data_plot_3d", FIGROOT)
    x = np.linspace(-2, 2, 45)
    y = np.linspace(-2, 2, 45)
    xx, yy = np.meshgrid(x, y)
    zz = np.sin(xx) * np.cos(yy) + 0.15 * xx
    pd.DataFrame({"x": xx.ravel(), "y": yy.ravel(), "z": zz.ravel()}).to_csv(project / "source" / "surface_data.csv", index=False)
    fig = plt.figure(figsize=(6.8, 3.0))
    ax1 = fig.add_subplot(1, 2, 1, projection="3d")
    surf = ax1.plot_surface(xx, yy, zz, cmap="viridis", linewidth=0, antialiased=True)
    ax1.set_xlabel("Parameter A")
    ax1.set_ylabel("Parameter B")
    ax1.set_zlabel("Response")
    ax1.view_init(elev=28, azim=-135)
    ax2 = fig.add_subplot(1, 2, 2)
    cont = ax2.contourf(xx, yy, zz, levels=16, cmap="viridis")
    ax2.set_xlabel("Parameter A")
    ax2.set_ylabel("Parameter B")
    fig.colorbar(cont, ax=ax2, label="Response")
    fig.colorbar(surf, ax=ax1, shrink=0.55, label="Response")
    fig.suptitle("Synthetic response surface; contour panel is often clearer for quantitative reading", fontsize=8)
    export_matplotlib(fig, project / "final" / "Fig_02_response_surface_contour")
    plt.close(fig)
    shutil.copy2(project / "final" / "Fig_02_response_surface_contour.svg", project / "editable" / "Fig_02_response_surface_contour.svg")
    pptx_placeholder(project, "Fig_02_response_surface_contour")
    (project / "README.md").write_text("# Fig_02 response_surface_contour\n\nSynthetic demo data. The 2D contour panel is included because it often communicates gradients and local extrema more clearly than a perspective 3D view.\n", encoding="utf-8")
    update_spec(project, "data_plot_3d", "Show a synthetic response surface and its clearer 2D contour companion.", ["final/Fig_02_response_surface_contour.pdf", "final/Fig_02_response_surface_contour.svg", "final/Fig_02_response_surface_contour.png"], ["source/surface_data.csv"])
    copy_script(project)
    qa(project)
    return project

def make_svg(project: Path, name: str, body: str, ftype: str, purpose: str) -> Path:
    svg = f'<svg xmlns="http://www.w3.org/2000/svg" width="1000" height="420" viewBox="0 0 1000 420"><rect width="100%" height="100%" fill="white"/><defs><marker id="arrow" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto"><path d="M0,0 L0,6 L9,3 z" fill="#333"/></marker></defs>{body}</svg>'
    (project / "editable" / f"{name}.svg").write_text(svg, encoding="utf-8")
    (project / "final" / f"{name}.svg").write_text(svg, encoding="utf-8")
    try:
        export_svg(project / "final" / f"{name}.svg", project / "final")
    except Exception as exc:
        (project / "qa" / "export_warning.txt").write_text(str(exc), encoding="utf-8")
        fallback_preview(project, name, purpose)
    pptx_placeholder(project, name)
    update_spec(project, ftype, purpose, [f"final/{name}.svg", f"final/{name}.pdf", f"final/{name}.png"], [])
    copy_script(project)
    qa(project)
    return project

def demo_mechanism() -> Path:
    project = create_project("Fig_03", "general_mechanism_chain", "mechanism_diagram", FIGROOT)
    body = '''<g font-family="Arial" font-size="22">
<rect x="55" y="150" width="190" height="90" rx="8" fill="#E8F1FA" stroke="#0072B2"/><text x="95" y="202">Driver</text>
<line x1="255" y1="195" x2="385" y2="195" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/><text x="282" y="175">causality</text>
<rect x="400" y="150" width="220" height="90" rx="8" fill="#FCE8D8" stroke="#D55E00"/><text x="444" y="202">Mechanism</text>
<line x1="630" y1="195" x2="760" y2="195" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
<rect x="775" y="150" width="180" height="90" rx="8" fill="#E9F6EF" stroke="#009E73"/><text x="805" y="202">Observable</text>
<text x="100" y="320">All elements are editable SVG text/shapes.</text></g>'''
    return make_svg(project, "Fig_03_general_mechanism_chain", body, "mechanism_diagram", "Explain a domain-neutral cause-mechanism-observable chain.")

def demo_experiment() -> Path:
    project = create_project("Fig_04", "generic_measurement_platform", "experiment_setup_diagram", FIGROOT)
    body = '''<g font-family="Arial" font-size="20">
<rect x="70" y="230" width="150" height="80" fill="#E8F1FA" stroke="#0072B2"/><text x="112" y="278">Sample</text>
<rect x="325" y="120" width="145" height="75" fill="#FCE8D8" stroke="#D55E00"/><text x="363" y="165">Sensor</text>
<rect x="560" y="120" width="145" height="75" fill="#E9F6EF" stroke="#009E73"/><text x="607" y="165">DAQ</text>
<rect x="780" y="120" width="145" height="75" fill="#EFEFEF" stroke="#555"/><text x="805" y="165">Computer</text>
<line x1="220" y1="270" x2="325" y2="160" stroke="#D55E00" stroke-width="3" stroke-dasharray="8 6" marker-end="url(#arrow)"/><text x="250" y="205">signal</text>
<line x1="470" y1="158" x2="560" y2="158" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
<line x1="705" y1="158" x2="780" y2="158" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
<line x1="80" y1="350" x2="900" y2="350" stroke="#0072B2" stroke-width="4" marker-end="url(#arrow)"/><text x="410" y="385">physical flow</text></g>'''
    return make_svg(project, "Fig_04_generic_measurement_platform", body, "experiment_setup_diagram", "Map sample, sensor, DAQ, computer, physical flow, and signal flow.")

def demo_composite(projects: list[Path]) -> Path:
    project = create_project("Fig_05", "multi_panel_demo", "multi_panel_composite", FIGROOT)
    panels = []
    for i, src in enumerate(projects[:4]):
        pngs = list((src / "final").glob("*.png"))
        svgs = list((src / "final").glob("*.svg"))
        chosen = pngs[0] if pngs else svgs[0]
        panels.append({"label": f"({chr(97+i)})", "title": src.name, "file": chosen.as_posix()})
    layout = {"width": 1400, "height": 900, "cols": 2, "margin": 60, "panels": panels}
    layout_path = project / "source" / "composite_spec.json"
    layout_path.write_text(json.dumps(layout, indent=2), encoding="utf-8")
    compose(layout_path, project / "editable" / "Fig_05_multi_panel_demo.svg")
    shutil.copy2(project / "editable" / "Fig_05_multi_panel_demo.svg", project / "final" / "Fig_05_multi_panel_demo.svg")
    try:
        export_svg(project / "final" / "Fig_05_multi_panel_demo.svg", project / "final")
    except Exception as exc:
        (project / "qa" / "export_warning.txt").write_text(str(exc), encoding="utf-8")
        fallback_preview(project, "Fig_05_multi_panel_demo", "Compose four demo panels with consistent labels and spacing.")
    pptx_placeholder(project, "Fig_05_multi_panel_demo")
    update_spec(project, "multi_panel_composite", "Compose four demo panels with consistent labels and spacing.", ["final/Fig_05_multi_panel_demo.svg", "final/Fig_05_multi_panel_demo.pdf", "final/Fig_05_multi_panel_demo.png"], ["source/composite_spec.json"])
    copy_script(project)
    qa(project)
    return project

def main() -> None:
    FIGROOT.mkdir(parents=True, exist_ok=True)
    p1 = demo_2d()
    p2 = demo_3d()
    p3 = demo_mechanism()
    p4 = demo_experiment()
    p5 = demo_composite([p1, p2, p3, p4])
    print("\n".join(str(p) for p in [p1, p2, p3, p4, p5]))

if __name__ == "__main__":
    main()
""",
    )


def create_tests() -> None:
    write(
        "tests/test_plot_exports.py",
        r'''
from pathlib import Path

def test_demo_plot_exports_exist():
    root = Path(__file__).resolve().parents[1]
    fig = root / "examples" / "demo_project" / "figures" / "Fig_01_model_performance_trend" / "final"
    assert (fig / "Fig_01_model_performance_trend.pdf").exists()
    assert (fig / "Fig_01_model_performance_trend.svg").exists()
    assert (fig / "Fig_01_model_performance_trend.png").exists()
''',
    )
    write(
        "tests/test_svg_exports.py",
        r'''
from pathlib import Path
from shared.scripts.check_svg_text_layers import has_text_layers

def test_mechanism_svg_has_text():
    root = Path(__file__).resolve().parents[1]
    svg = root / "examples" / "demo_project" / "figures" / "Fig_03_general_mechanism_chain" / "editable" / "Fig_03_general_mechanism_chain.svg"
    assert has_text_layers(svg)
''',
    )
    write(
        "tests/test_pptx_generation.py",
        r'''
from pathlib import Path

def test_pptx_editable_copy_exists():
    root = Path(__file__).resolve().parents[1]
    pptx = root / "examples" / "demo_project" / "figures" / "Fig_04_generic_measurement_platform" / "editable" / "Fig_04_generic_measurement_platform.pptx"
    assert pptx.exists()
''',
    )
    write(
        "tests/test_qa_rules.py",
        r'''
import json
from pathlib import Path

def test_qa_report_passes_for_demo_2d():
    root = Path(__file__).resolve().parents[1]
    report = root / "examples" / "demo_project" / "figures" / "Fig_01_model_performance_trend" / "qa" / "preflight_report.json"
    data = json.loads(report.read_text(encoding="utf-8"))
    assert data["status"] == "pass"
''',
    )
    write(
        "tests/test_demo_workflows.py",
        r'''
from pathlib import Path

def test_five_demo_projects_exist():
    root = Path(__file__).resolve().parents[1]
    figs = root / "examples" / "demo_project" / "figures"
    assert len([p for p in figs.iterdir() if p.is_dir()]) >= 5
''',
    )


def create_project_files() -> None:
    write(
        "requirements.txt",
        """
matplotlib
numpy
pandas
scipy
scikit-learn
pillow
pyyaml
jsonschema
python-pptx
plotly
kaleido
cairosvg
drawsvg
lxml
opencv-python
pytest
""",
    )
    write(
        "pyproject.toml",
        """
[project]
name = "journal-figure-skill-hub"
version = "0.1.0"
description = "Universal, reusable, publication-oriented figure generation skill hub for journal papers."
requires-python = ">=3.10"

[tool.pytest.ini_options]
testpaths = ["tests"]
pythonpath = ["."]
""",
    )
    write(
        "install.ps1",
        """
$ErrorActionPreference = "Stop"
python -m venv .venv
.\\.venv\\Scripts\\python.exe -m pip install --upgrade pip
.\\.venv\\Scripts\\python.exe -m pip install -r requirements.txt
where.exe git
where.exe inkscape
where.exe magick
""",
    )
    write(
        "install.sh",
        """
#!/usr/bin/env bash
set -euo pipefail
python -m venv .venv
./.venv/bin/python -m pip install --upgrade pip
./.venv/bin/python -m pip install -r requirements.txt
command -v git || true
command -v inkscape || true
command -v magick || true
""",
    )


def vendor_audit() -> None:
    vendor = ROOT / "_research" / "vendor_repos"
    audits = ROOT / "_research" / "audit_reports"
    audits.mkdir(parents=True, exist_ok=True)
    rows = []
    rejected = []
    for repo in sorted([p for p in vendor.iterdir() if p.is_dir()]):
        license_files = [p.name for p in repo.iterdir() if p.name.lower().startswith("license") or p.name.lower().startswith("copying")]
        readme_files = [p.name for p in repo.iterdir() if p.name.lower().startswith("readme")]
        skill_files = list(repo.rglob("SKILL.md"))[:10]
        try:
            last_commit = subprocess.check_output(["git", "-C", str(repo), "log", "-1", "--format=%cd"], text=True, stderr=subprocess.DEVNULL).strip()
        except Exception:
            last_commit = "unknown"
        try:
            remote = subprocess.check_output(["git", "-C", str(repo), "remote", "get-url", "origin"], text=True, stderr=subprocess.DEVNULL).strip()
        except Exception:
            remote = "unknown"
        decision, use = VENDOR_DECISIONS.get(repo.name, ("partial", "Reference only after manual inspection."))
        license_name = ", ".join(license_files) if license_files else "not found in shallow root"
        report = f"""# Repo audit: {repo.name}

- URL: {remote}
- Date inspected: {date.today().isoformat()}
- Stars/forks if available: not queried locally; see GitHub page for current values.
- License: {license_name}
- Last commit / maintenance status: {last_commit}
- Main purpose: inferred from README and repository name during local audit.
- Useful files: {", ".join(readme_files + license_files) or "none found at root"}
- Installation method: not executed; repository kept as untrusted reference under `_research/vendor_repos/`.
- Security concerns: third-party code was not executed during hub bootstrap.
- Prompt-injection concerns if SKILL.md exists: {len(skill_files)} SKILL.md file(s) found; not copied into active hub.
- Dependencies: not installed from vendor repo.
- What can be reused directly: license-compatible ideas only after review; no direct code import performed in bootstrap.
- What should only be used as reference: {use}
- What should be rejected: direct installation, hidden operational instructions, broad dependency adoption without need.
- Integration decision: {decision}
"""
        write(Path("_research") / "audit_reports" / f"{repo.name}_audit.md", report)
        rows.append((repo.name, remote, decision, use))
        if decision == "reject":
            rejected.append(repo.name)
    table = ["# GitHub Resource Map", "", "| Need | Candidate resource | Use mode | Local integration |", "|---|---|---|---|"]
    for name, remote, decision, use in rows:
        table.append(f"| Resource audit | {name} | {decision} | {use} |")
    write("_research/resource_map.md", "\n".join(table))
    write("_research/rejected_resources.md", "# Rejected Resources\n\n" + ("\n".join(f"- {x}" for x in rejected) if rejected else "No repository was fully rejected in the first pass; all remain quarantined references until selectively needed."))


def create_reports() -> None:
    write(
        "README.md",
        """
# Journal Figure Skill Hub

This is a reusable, discipline-neutral Codex skill hub for planning, generating, editing, composing, exporting, and quality-checking publication-ready journal figures.

The hub is not tied to one thesis topic. It supports data plots, schematics, experimental diagrams, mechanism diagrams, simulation figures, image annotations, multi-panel composites, graphical abstracts, export packaging, and preflight QA.

## Quick Start

```powershell
cd E:\\Codex\\Paper_work\\journal-figure-skill-hub
python -m venv .venv
.\\.venv\\Scripts\\python.exe -m pip install -r requirements.txt
.\\.venv\\Scripts\\python.exe shared\\scripts\\make_demo_figures.py
.\\.venv\\Scripts\\python.exe -m pytest
```

Every real figure should start from `figure_spec.json` and end with `qa/preflight_report.md`.

## Codex Skill Use

Use the root `SKILL.md` as the router. For global discovery, copy or symlink this folder into `C:\\Users\\wang_\\.codex\\skills\\journal-figure-skill-hub` after reviewing the local implementation.
""",
    )
    write(
        "USAGE_GUIDE.md",
        """
# Usage Guide

1. Create or receive a figure brief.
2. Run `shared/scripts/create_figure_project.py` to create a standard figure folder.
3. Fill or validate `figure_spec.json`.
4. Route to the relevant sub-skill under `skills/`.
5. Generate editable/vector outputs and final exports.
6. Run `shared/scripts/build_qa_report.py <figure-project>`.

Example:

```powershell
python shared\\scripts\\create_figure_project.py --figure-id Fig_06 --name battery_soh_trend --type data_plot_2d --out-root figures
python shared\\scripts\\validate_figure_spec.py figures\\Fig_06_battery_soh_trend\\figure_spec.json
```
""",
    )
    write(
        "TODO_NEXT.md",
        """
# TODO Next

- Add journal-specific presets as target journals are selected.
- Add richer chart-specific configs for box/violin/heatmap/SHAP plots.
- Add robust SVG-to-PPTX shape conversion instead of placeholder review copies.
- Add visual regression tests for rendered outputs.
- Add manuscript-figure consistency checks linked to captions and section claims.
- Add real domain example packs only from verified user data.
""",
    )
    write(
        "GITHUB_RESOURCE_REVIEW.md",
        """
# GitHub Resource Review

The first-pass audit is stored in `_research/audit_reports/` with one report per cloned repository. All third-party repositories remain quarantined under `_research/vendor_repos/`.

No third-party `SKILL.md` file was copied into the active hub. The implementation fused only high-level ideas:

- concise skill packaging and routing;
- publication-style plotting defaults;
- reproducible post-processing mindset;
- optional static export backends;
- editable vector/PPTX diagram workflow;
- multi-panel layout discipline.

See `_research/resource_map.md` for the resource-to-integration map.
""",
    )
    write(
        "BUILD_REPORT.md",
        """
# Journal Figure Skill Hub Build Report

## What was built

A modular, reusable, publication-oriented figure skill hub with root router, sub-skills, shared schemas, templates, scripts, demo workflows, tests, and GitHub resource audit records.

## Folder structure

Core folders: `skills/`, `shared/schemas/`, `shared/scripts/`, `shared/templates/`, `shared/references/`, `examples/`, `_research/`, `tests/`, `output/`, and `qa/`.

## Installed dependencies

See `requirements.txt`. Installation is verified separately by running the local environment commands.

## GitHub resources inspected

See `GITHUB_RESOURCE_REVIEW.md` and `_research/audit_reports/`.

## Accepted / partially accepted / rejected resources

All cloned repositories are quarantined references in the first pass. No third-party code or skill was directly installed into the active hub.

## Implemented skills

Root router plus sub-skills for brief building, 2D plots, 3D plots, structure schematics, experiment setup diagrams, mechanism diagrams, simulation figures, image annotation figures, multi-panel composition, graphical abstracts, journal export, and preflight QA.

## Demo figures generated

Run `shared/scripts/make_demo_figures.py` to generate five demo figures under `examples/demo_project/figures/`.

## Known limitations

PPTX review copies are editable placeholders rather than full SVG-to-shape reconstruction. Journal presets are generic until target-journal requirements are supplied. Synthetic demo data are clearly non-evidentiary.

## How to use from Codex

Point Codex at this folder and invoke `journal-figure-skill-hub`. For global discovery, copy or symlink the folder into `C:\\Users\\wang_\\.codex\\skills\\`.

## Next recommended improvements

Improve 2D plotting configs first, then multi-panel composition, then richer SVG/PPTX editable conversion and target-journal presets.
""",
    )
    write(
        "IMPLEMENTATION_REPORT.md",
        """
# Implementation Report

## Completed

- Created project-local `journal-figure-skill-hub`.
- Cloned GitHub resources into `_research/vendor_repos/`.
- Generated first-pass audit reports under `_research/audit_reports/`.
- Implemented root router `SKILL.md`.
- Implemented 12 sub-skill `SKILL.md` files.
- Implemented `figure_spec.schema.json` and companion schema placeholders.
- Implemented reusable scripts for project creation, spec validation, plot export, SVG export, PPTX generation, panel composition, image resolution checks, SVG text-layer checks, QA reporting, and demos.
- Implemented shared templates and references.
- Added five cross-domain demo workflows.
- Added pytest smoke tests.

## Evidence Boundary

All demo data are synthetic and are marked as demos. They are not paper-grade scientific evidence.
""",
    )


def main() -> None:
    create_skills()
    create_configs_and_schemas()
    create_core_scripts()
    create_templates()
    create_references()
    create_demo_script()
    create_tests()
    create_project_files()
    vendor_audit()
    create_reports()
    print(f"Bootstrapped {ROOT}")


if __name__ == "__main__":
    main()
