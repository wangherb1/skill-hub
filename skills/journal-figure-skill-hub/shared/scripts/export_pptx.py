from __future__ import annotations
import argparse, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

def build_pptx(shape_spec: Path, output: Path) -> Path:
    spec = json.loads(shape_spec.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(spec.get("width_in", 10))
    prs.slide_height = Inches(spec.get("height_in", 5.625))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for s in spec.get("shapes", []):
        x, y, w, h = [Inches(float(s.get(k, 0))) for k in ["x", "y", "w", "h"]]
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.text = s.get("text", "")
        if shape.text_frame and shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].font.size = Pt(float(s.get("font_pt", 12)))
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output

def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("shape_spec")
    p.add_argument("output")
    args = p.parse_args()
    print(build_pptx(Path(args.shape_spec), Path(args.output)))

if __name__ == "__main__":
    main()
