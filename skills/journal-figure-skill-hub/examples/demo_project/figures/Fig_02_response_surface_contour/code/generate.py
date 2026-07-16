from __future__ import annotations

import json
import math
import shutil
import subprocess
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from create_figure_project import create_project
from export_plot import export_matplotlib
from export_svg import export_svg
from build_qa_report import qa
from compose_panels import compose

ROOT = Path(__file__).resolve().parents[2]
FIGROOT = ROOT / "examples" / "demo_project" / "figures"

def update_spec(project: Path, figure_type: str, purpose: str, outputs: list[str], sources: list[str]) -> None:
    spec = json.loads((project / "figure_spec.json").read_text(encoding="utf-8"))
    spec["figure_type"] = figure_type
    spec["scientific_purpose"] = purpose
    spec["reader_takeaway_5s"] = purpose
    spec["data_sources"] = sources
    spec["panels"][0]["message"] = purpose
    spec["panels"][0]["output_files"] = outputs
    (project / "figure_spec.json").write_text(json.dumps(spec, indent=2, ensure_ascii=False), encoding="utf-8")

def copy_script(project: Path) -> None:
    shutil.copy2(__file__, project / "code" / "generate.py")

def pptx_placeholder(project: Path, name: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.5), Inches(8.8), Inches(4.5))
    box.text = f"Editable review copy: {name}"
    box.text_frame.paragraphs[0].font.size = Pt(24)
    prs.save(project / "editable" / f"{name}.pptx")

def fallback_preview(project: Path, name: str, title: str) -> None:
    fig, ax = plt.subplots(figsize=(6.0, 2.6))
    ax.axis("off")
    ax.text(0.5, 0.72, title, ha="center", va="center", fontsize=12, fontweight="bold")
    ax.text(0.5, 0.48, "Editable vector source is the SVG file in editable/ and final/.", ha="center", va="center", fontsize=9)
    ax.text(0.5, 0.30, "This PDF/PNG is a renderer fallback preview for environments without Cairo/Inkscape.", ha="center", va="center", fontsize=8)
    export_matplotlib(fig, project / "final" / name, dpi=600)
    plt.close(fig)

def demo_2d() -> Path:
    project = create_project("Fig_01", "model_performance_trend", "data_plot_2d", FIGROOT)
    x = np.linspace(0, 10, 28)
    df = pd.DataFrame({"x": x, "baseline": 0.58 + 0.18 * (1 - np.exp(-x / 4)), "proposed": 0.62 + 0.27 * (1 - np.exp(-x / 3))})
    df.to_csv(project / "source" / "raw_data.csv", index=False)
    fig, ax = plt.subplots(figsize=(3.6, 2.4))
    ax.plot(df["x"], df["baseline"], "o-", label="Baseline")
    ax.plot(df["x"], df["proposed"], "s-", label="Proposed")
    ax.set_xlabel("Training size (a.u.)")
    ax.set_ylabel("Score (a.u.)")
    ax.legend(frameon=False)
    ax.grid(alpha=0.25)
    export_matplotlib(fig, project / "final" / "Fig_01_model_performance_trend")
    plt.close(fig)
    shutil.copy2(project / "final" / "Fig_01_model_performance_trend.svg", project / "editable" / "Fig_01_model_performance_trend.svg")
    pptx_placeholder(project, "Fig_01_model_performance_trend")
    update_spec(project, "data_plot_2d", "Compare synthetic model-performance trends with reproducible source data.", ["final/Fig_01_model_performance_trend.pdf", "final/Fig_01_model_performance_trend.svg", "final/Fig_01_model_performance_trend.png"], ["source/raw_data.csv"])
    copy_script(project)
    qa(project)
    return project

def demo_3d() -> Path:
    project = create_project("Fig_02", "response_surface_contour", "data_plot_3d", FIGROOT)
    x = np.linspace(-2, 2, 45)
    y = np.linspace(-2, 2, 45)
    xx, yy = np.meshgrid(x, y)
    zz = np.sin(xx) * np.cos(yy) + 0.15 * xx
    pd.DataFrame({"x": xx.ravel(), "y": yy.ravel(), "z": zz.ravel()}).to_csv(project / "source" / "surface_data.csv", index=False)
    fig = plt.figure(figsize=(6.8, 3.0))
    ax1 = fig.add_subplot(1, 2, 1, projection="3d")
    surf = ax1.plot_surface(xx, yy, zz, cmap="viridis", linewidth=0, antialiased=True)
    ax1.set_xlabel("Parameter A")
    ax1.set_ylabel("Parameter B")
    ax1.set_zlabel("Response")
    ax1.view_init(elev=28, azim=-135)
    ax2 = fig.add_subplot(1, 2, 2)
    cont = ax2.contourf(xx, yy, zz, levels=16, cmap="viridis")
    ax2.set_xlabel("Parameter A")
    ax2.set_ylabel("Parameter B")
    fig.colorbar(cont, ax=ax2, label="Response")
    fig.colorbar(surf, ax=ax1, shrink=0.55, label="Response")
    fig.suptitle("Synthetic response surface; contour panel is often clearer for quantitative reading", fontsize=8)
    export_matplotlib(fig, project / "final" / "Fig_02_response_surface_contour")
    plt.close(fig)
    shutil.copy2(project / "final" / "Fig_02_response_surface_contour.svg", project / "editable" / "Fig_02_response_surface_contour.svg")
    pptx_placeholder(project, "Fig_02_response_surface_contour")
    (project / "README.md").write_text("# Fig_02 response_surface_contour\n\nSynthetic demo data. The 2D contour panel is included because it often communicates gradients and local extrema more clearly than a perspective 3D view.\n", encoding="utf-8")
    update_spec(project, "data_plot_3d", "Show a synthetic response surface and its clearer 2D contour companion.", ["final/Fig_02_response_surface_contour.pdf", "final/Fig_02_response_surface_contour.svg", "final/Fig_02_response_surface_contour.png"], ["source/surface_data.csv"])
    copy_script(project)
    qa(project)
    return project

def make_svg(project: Path, name: str, body: str, ftype: str, purpose: str) -> Path:
    svg = f'<svg xmlns="http://www.w3.org/2000/svg" width="1000" height="420" viewBox="0 0 1000 420"><rect width="100%" height="100%" fill="white"/><defs><marker id="arrow" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto"><path d="M0,0 L0,6 L9,3 z" fill="#333"/></marker></defs>{body}</svg>'
    (project / "editable" / f"{name}.svg").write_text(svg, encoding="utf-8")
    (project / "final" / f"{name}.svg").write_text(svg, encoding="utf-8")
    try:
        export_svg(project / "final" / f"{name}.svg", project / "final")
    except Exception as exc:
        (project / "qa" / "export_warning.txt").write_text(str(exc), encoding="utf-8")
        fallback_preview(project, name, purpose)
    pptx_placeholder(project, name)
    update_spec(project, ftype, purpose, [f"final/{name}.svg", f"final/{name}.pdf", f"final/{name}.png"], [])
    copy_script(project)
    qa(project)
    return project

def demo_mechanism() -> Path:
    project = create_project("Fig_03", "general_mechanism_chain", "mechanism_diagram", FIGROOT)
    body = '''<g font-family="Arial" font-size="22">
<rect x="55" y="150" width="190" height="90" rx="8" fill="#E8F1FA" stroke="#0072B2"/><text x="95" y="202">Driver</text>
<line x1="255" y1="195" x2="385" y2="195" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/><text x="282" y="175">causality</text>
<rect x="400" y="150" width="220" height="90" rx="8" fill="#FCE8D8" stroke="#D55E00"/><text x="444" y="202">Mechanism</text>
<line x1="630" y1="195" x2="760" y2="195" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
<rect x="775" y="150" width="180" height="90" rx="8" fill="#E9F6EF" stroke="#009E73"/><text x="805" y="202">Observable</text>
<text x="100" y="320">All elements are editable SVG text/shapes.</text></g>'''
    return make_svg(project, "Fig_03_general_mechanism_chain", body, "mechanism_diagram", "Explain a domain-neutral cause-mechanism-observable chain.")

def demo_experiment() -> Path:
    project = create_project("Fig_04", "generic_measurement_platform", "experiment_setup_diagram", FIGROOT)
    body = '''<g font-family="Arial" font-size="20">
<rect x="70" y="230" width="150" height="80" fill="#E8F1FA" stroke="#0072B2"/><text x="112" y="278">Sample</text>
<rect x="325" y="120" width="145" height="75" fill="#FCE8D8" stroke="#D55E00"/><text x="363" y="165">Sensor</text>
<rect x="560" y="120" width="145" height="75" fill="#E9F6EF" stroke="#009E73"/><text x="607" y="165">DAQ</text>
<rect x="780" y="120" width="145" height="75" fill="#EFEFEF" stroke="#555"/><text x="805" y="165">Computer</text>
<line x1="220" y1="270" x2="325" y2="160" stroke="#D55E00" stroke-width="3" stroke-dasharray="8 6" marker-end="url(#arrow)"/><text x="250" y="205">signal</text>
<line x1="470" y1="158" x2="560" y2="158" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
<line x1="705" y1="158" x2="780" y2="158" stroke="#333" stroke-width="3" marker-end="url(#arrow)"/>
<line x1="80" y1="350" x2="900" y2="350" stroke="#0072B2" stroke-width="4" marker-end="url(#arrow)"/><text x="410" y="385">physical flow</text></g>'''
    return make_svg(project, "Fig_04_generic_measurement_platform", body, "experiment_setup_diagram", "Map sample, sensor, DAQ, computer, physical flow, and signal flow.")

def demo_composite(projects: list[Path]) -> Path:
    project = create_project("Fig_05", "multi_panel_demo", "multi_panel_composite", FIGROOT)
    panels = []
    for i, src in enumerate(projects[:4]):
        pngs = list((src / "final").glob("*.png"))
        svgs = list((src / "final").glob("*.svg"))
        chosen = pngs[0] if pngs else svgs[0]
        panels.append({"label": f"({chr(97+i)})", "title": src.name, "file": chosen.as_posix()})
    layout = {"width": 1400, "height": 900, "cols": 2, "margin": 60, "panels": panels}
    layout_path = project / "source" / "composite_spec.json"
    layout_path.write_text(json.dumps(layout, indent=2), encoding="utf-8")
    compose(layout_path, project / "editable" / "Fig_05_multi_panel_demo.svg")
    shutil.copy2(project / "editable" / "Fig_05_multi_panel_demo.svg", project / "final" / "Fig_05_multi_panel_demo.svg")
    try:
        export_svg(project / "final" / "Fig_05_multi_panel_demo.svg", project / "final")
    except Exception as exc:
        (project / "qa" / "export_warning.txt").write_text(str(exc), encoding="utf-8")
        fallback_preview(project, "Fig_05_multi_panel_demo", "Compose four demo panels with consistent labels and spacing.")
    pptx_placeholder(project, "Fig_05_multi_panel_demo")
    update_spec(project, "multi_panel_composite", "Compose four demo panels with consistent labels and spacing.", ["final/Fig_05_multi_panel_demo.svg", "final/Fig_05_multi_panel_demo.pdf", "final/Fig_05_multi_panel_demo.png"], ["source/composite_spec.json"])
    copy_script(project)
    qa(project)
    return project

def main() -> None:
    FIGROOT.mkdir(parents=True, exist_ok=True)
    p1 = demo_2d()
    p2 = demo_3d()
    p3 = demo_mechanism()
    p4 = demo_experiment()
    p5 = demo_composite([p1, p2, p3, p4])
    print("\n".join(str(p) for p in [p1, p2, p3, p4, p5]))

if __name__ == "__main__":
    main()
